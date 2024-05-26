import streamlit as st
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings  # Updated import
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFaceHub
import os
from dotenv import load_dotenv
import tempfile
from transformers import pipeline
import pandas as pd
import io

# Load environment variables from .env file
load_dotenv()
huggingfacehub_api_token = os.getenv("HUGGINGFACEHUB_API_TOKEN")

# Check if the API key is loaded correctly
if huggingfacehub_api_token is None:
    raise ValueError("No Hugging Face API token found. Please set the HUGGINGFACEHUB_API_TOKEN environment variable in your .env file.")

# Extract text from a PDF file
def get_pdf_text(pdf_file):
    text = ""
    pdf_reader = PdfReader(pdf_file)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

# Extract text from a DOCX file
def get_word_text(docx_file):
    document = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in document.paragraphs])
    return text

# Extract text from a TXT file
def read_text_file(txt_file):
    text = txt_file.getvalue().decode('utf-8')
    return text

# Extract text from a PPTX file
def get_pptx_text(pptx_file):
    presentation = Presentation(pptx_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Handle CSV file input
def handle_csv_file(csv_file, user_question):
    # Read the CSV file
    csv_text = csv_file.read().decode("utf-8")
    
    # Create a DataFrame from the CSV text
    df = pd.read_csv(io.StringIO(csv_text))
    df = df.astype(str)
    
    # Initialize a Hugging Face table-question-answering pipeline
    qa_pipeline = pipeline("table-question-answering", model="google/tapas-large-finetuned-wtq")
    
    # Use the pipeline to answer the question
    response = qa_pipeline(table=df, query=user_question)
    
    # Display the answer
    st.write(response['answer'])

def combine_text(text_list):
    return "\n".join(text_list)

# Split text into chunks
def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

# Create vector store from text chunks
def get_vectorstore(text_chunks):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")  # Updated to use HuggingFaceEmbeddings
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

# Create conversation chain
def get_conversation_chain(vectorstore):
    llm = HuggingFaceHub(
        repo_id="google/flan-t5-xxl",
        huggingfacehub_api_token=huggingfacehub_api_token,  # Pass the token directly
        model_kwargs={"temperature": 0.5, "max_length": 200}  # Adjusted max_length to stay within token limit
    )
    memory = ConversationBufferMemory(
        memory_key='chat_history', 
        return_messages=True
    )
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

# Handle user input
def handle_userinput(user_question):
    if st.session_state.conversation is not None:
        response = st.session_state.conversation({'question': user_question})
        st.session_state.chat_history = response['chat_history']

        for i, message in enumerate(st.session_state.chat_history):
            if i % 2 == 0:
                st.write(user_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
            else:
                st.write(bot_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
    else:
        # Handle the case when conversation is not initialized
        st.write("Please upload and process your documents first.")

# Main function
def main():
    load_dotenv()
    st.set_page_config(
        page_title="File Chatbot",
        page_icon=":books:",
        layout="wide"
    )
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Chat with your multiple files:")
    user_question = st.text_input("Ask a question about your documents:")

    # Initialize variables to hold uploaded files
    csv_file = None
    other_files = []

    with st.sidebar:
        st.subheader("Your documents")
        files = st.file_uploader("Upload your files here and click on 'Process'", accept_multiple_files=True)
        
        for file in files:
            if file.name.lower().endswith('.csv'):
                csv_file = file  # Store the CSV file
            else:
                other_files.append(file)  # Store other file types

        # Initialize empty lists for each file type
        pdf_texts = []
        word_texts = []
        txt_texts = []
        pptx_texts = []

        if st.button("Process"):
            with st.spinner("Processing"):
                for file in other_files:
                    if file.name.lower().endswith('.pdf'):
                        pdf_texts.append(get_pdf_text(file))
                    elif file.name.lower().endswith('.docx'):
                        word_texts.append(get_word_text(file))
                    elif file.name.lower().endswith('.txt'):
                        txt_texts.append(read_text_file(file))
                    elif file.name.lower().endswith('.pptx'):
                        pptx_texts.append(get_pptx_text(file))

                # Combine text from different file types
                combined_text = combine_text(pdf_texts + word_texts + txt_texts + pptx_texts)

                # Split the combined text into chunks
                text_chunks = get_text_chunks(combined_text)

                # Create vector store and conversation chain if non-CSV documents are uploaded
                if len(other_files) > 0:
                    vectorstore = get_vectorstore(text_chunks)
                    st.session_state.conversation = get_conversation_chain(vectorstore)
                else:
                    vectorstore = None  # No need for vectorstore with CSV file

    # Handle user input for CSV file separately
    if csv_file is not None and user_question:
        handle_csv_file(csv_file, user_question)
    
    # Handle user input for text-based files
    if user_question:
        handle_userinput(user_question)

if __name__ == '__main__':
    main()
